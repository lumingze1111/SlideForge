"""
PPTX Exporter - Convert slide outline + color scheme to .pptx using python-pptx.

Supports gradient backgrounds and speaker notes.
"""

import re
import math
from pathlib import Path
from typing import Optional, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from slideforge.agents.html_generator import PresentationOutline, SlideContent


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN_H = Inches(0.8)
MARGIN_V = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN_H * 2


def _parse_hex(color: str) -> Optional[RGBColor]:
    m = re.search(r'#([0-9a-fA-F]{6})', color)
    if m:
        h = m.group(1)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return None


def _parse_all_hex(color: str) -> List[str]:
    """Extract all hex colors from a gradient string, in order."""
    return re.findall(r'#([0-9a-fA-F]{6})', color)


def _parse_gradient_stops(color_str: str) -> List[tuple]:
    """Parse CSS gradient color stops with positions.

    Returns list of (hex_color, position_0_to_1000) tuples.
    Handles: no positions, % positions, or mixed.
    """
    # Match each color stop: #rrggbb with optional position
    pattern = r'#([0-9a-fA-F]{6})\s*(?:(\d+(?:\.\d+)?)\s*%)?'
    stops = re.findall(pattern, color_str)

    if not stops:
        return []

    result = []
    explicit_positions = [s for s in stops if s[1]]

    if explicit_positions:
        # Use explicit positions from CSS
        for hex_color, pos_pct in stops:
            if pos_pct:
                pos = int(float(pos_pct) * 10)  # % → 0-1000 scale
            else:
                pos = None
            result.append((hex_color, pos))

        # Fill in missing positions with linear interpolation
        for i in range(len(result)):
            if result[i][1] is None:
                # Find nearest left and right known positions
                left_idx = right_idx = None
                for j in range(i - 1, -1, -1):
                    if result[j][1] is not None:
                        left_idx = j
                        break
                for j in range(i + 1, len(result)):
                    if result[j][1] is not None:
                        right_idx = j
                        break

                if left_idx is not None and right_idx is not None:
                    # Interpolate between known positions
                    left_pos = result[left_idx][1]
                    right_pos = result[right_idx][1]
                    frac = (i - left_idx) / (right_idx - left_idx)
                    pos = int(left_pos + frac * (right_pos - left_pos))
                elif left_idx is not None:
                    pos = result[left_idx][1]
                elif right_idx is not None:
                    pos = result[right_idx][1]
                else:
                    pos = 0

                result[i] = (result[i][0], pos)
    else:
        # No explicit positions: evenly distribute 0 to 1000
        n = len(stops)
        for i, (hex_color, _) in enumerate(stops):
            pos = int(i * 1000 / (n - 1))
            result.append((hex_color, pos))

    return result


def _parse_gradient_angle(color_str: str) -> int:
    """Extract gradient angle from CSS gradient, default 0."""
    m = re.search(r'linear-gradient\((\d+)deg', color_str)
    if m:
        return int(m.group(1))
    return 0


def _apply_gradient_bg(slide, color_str: str) -> None:
    """Apply a real gradient background to the slide using XML manipulation."""
    bg = slide.background
    fill_xml = bg._element.find(qn('p:bgPr'))
    if fill_xml is None:
        # Create bgPr element if it doesn't exist (e.g. on blank layouts)
        fill_xml = etree.SubElement(bg._element, qn('p:bgPr'))

    stops = _parse_gradient_stops(color_str)
    if not stops or len(stops) < 2:
        # Not a gradient or single color - fallback to solid
        rgb = _parse_hex(color_str)
        if rgb:
            bg.fill.solid()
            bg.fill.fore_color.rgb = rgb
        return

    angle = _parse_gradient_angle(color_str) if 'linear-gradient' in color_str else 0
    # CSS 0deg = bottom→top, OOXML 0° = left→right
    # Conversion: OOXML_angle = (CSS_angle - 90) % 360
    ooxml_angle = (angle - 90) % 360

    # Build gradient XML
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    # Remove existing fill if any
    for child in list(fill_xml):
        fill_xml.remove(child)

    grad_fill = etree.SubElement(fill_xml, qn('a:gradFill'))
    grad_fill.set('rotWithShape', '1')

    gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))
    for hex_color, pos in stops:
        gs = etree.SubElement(gs_lst, qn('a:gs'))
        gs.set('pos', str(pos))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', hex_color.upper())

    lin = etree.SubElement(grad_fill, qn('a:lin'))
    lin.set('ang', str(ooxml_angle))
    lin.set('scaled', '1')


def _add_text_box(
    slide, text: str, left: Emu, top: Emu, width: Emu, height: Emu,
    font_size: int, bold: bool = False, color: str = "#ffffff",
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    rgb = _parse_hex(color)
    if rgb:
        run.font.color.rgb = rgb


def _add_rect(
    slide, left: Emu, top: Emu, width: Emu, height: Emu,
    fill_color: str, line_color: Optional[str] = None,
    shape_type: int = 1, line_width: float = 1,
) -> None:
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    fill = shape.fill
    rgb = _parse_hex(fill_color)
    if rgb:
        fill.solid()
        fill.fore_color.rgb = rgb
    else:
        fill.background()
    if line_color:
        line_rgb = _parse_hex(line_color)
        if line_rgb:
            shape.line.color.rgb = line_rgb
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()


def _add_notes(slide, notes_text: str) -> None:
    """Add speaker notes to a slide."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = notes_text


# ─── Slide renderers ────────────────────────────────────────────────────────


def _render_cover(slide, s: SlideContent, colors: dict) -> None:
    bg = colors.get("gradient_bg", colors.get("background", "#1a1a2e"))
    primary = colors.get("primary", "#7c3aed")
    accent = colors.get("accent", "#f59e0b")
    text_sec = colors.get("text_secondary", "#94a3b8")

    _apply_gradient_bg(slide, bg)
    # 80px × 5px accent line → 0.83" × 0.05" at 96dpi
    _add_rect(slide, SLIDE_W / 2 - Inches(0.42), Inches(1.5), Inches(0.83), Inches(0.05), accent)
    _add_text_box(slide, s.title, MARGIN_H, Inches(2.0), CONTENT_W, Inches(2.0),
                  font_size=42, bold=True, color=primary, align=PP_ALIGN.CENTER)
    if s.subtitle:
        _add_text_box(slide, s.subtitle, MARGIN_H, Inches(4.2), CONTENT_W, Inches(1.2),
                      font_size=17, color=text_sec, align=PP_ALIGN.CENTER)
    _add_rect(slide, SLIDE_W / 2 - Inches(0.42), Inches(5.8), Inches(0.83), Inches(0.05), accent)
    _add_notes(slide, s.notes)


def _render_section(slide, s: SlideContent, colors: dict) -> None:
    bg = colors.get("gradient_bg", colors.get("background", "#1a1a2e"))
    primary = colors.get("primary", "#7c3aed")
    accent = colors.get("accent", "#f59e0b")
    text_sec = colors.get("text_secondary", "#94a3b8")

    _apply_gradient_bg(slide, bg)
    # 60px × 4px accent line → 0.63" × 0.04" at 96dpi
    _add_rect(slide, MARGIN_H, Inches(2.8), Inches(0.63), Inches(0.04), accent)
    _add_text_box(slide, s.title, MARGIN_H, Inches(2.0), CONTENT_W, Inches(1.6),
                  font_size=36, bold=True, color=primary)
    if s.subtitle:
        _add_text_box(slide, s.subtitle, MARGIN_H, Inches(3.2), CONTENT_W, Inches(1.0),
                      font_size=15, color=text_sec)
    _add_notes(slide, s.notes)


def _render_content(slide, s: SlideContent, colors: dict) -> None:
    bg = colors.get("gradient_bg", colors.get("background", "#1a1a2e"))
    primary = colors.get("primary", "#7c3aed")
    accent = colors.get("accent", "#f59e0b")
    text_sec = colors.get("text_secondary", "#94a3b8")

    _apply_gradient_bg(slide, bg)
    _add_text_box(slide, s.title, MARGIN_H, MARGIN_V, CONTENT_W, Inches(0.9),
                  font_size=29, bold=True, color=primary)
    # 80px × 3px accent line → 0.83" × 0.03" at 96dpi
    _add_rect(slide, MARGIN_H, Inches(1.6), Inches(0.83), Inches(0.03), accent)
    y = Inches(1.9)
    for bullet in s.bullets:
        _add_text_box(slide, "▸", MARGIN_H, y, Inches(0.3), Inches(0.4), 14, color=accent)
        _add_text_box(slide, bullet, MARGIN_H + Inches(0.3), y, CONTENT_W - Inches(0.3), Inches(0.5),
                      font_size=14, color=text_sec)
        y += Inches(0.5)
    _add_notes(slide, s.notes)


def _render_two_column(slide, s: SlideContent, colors: dict) -> None:
    bg = colors.get("gradient_bg", colors.get("background", "#1a1a2e"))
    primary = colors.get("primary", "#7c3aed")
    accent = colors.get("accent", "#f59e0b")
    text_sec = colors.get("text_secondary", "#94a3b8")
    surface = colors.get("surface", colors.get("card_bg", "#1e293b"))
    border = colors.get("border", "#475569")

    _apply_gradient_bg(slide, bg)
    _add_text_box(slide, s.title, MARGIN_H, MARGIN_V, CONTENT_W, Inches(0.9), 27, True, primary)
    # 80px × 3px accent line
    _add_rect(slide, MARGIN_H, Inches(1.6), Inches(0.83), Inches(0.03), accent)

    col_w = (CONTENT_W - Inches(0.3)) / 2
    mid = len(s.bullets) // 2
    left_bullets = s.bullets[:mid]
    right_bullets = s.bullets[mid:]

    # ROUNDED_RECTANGLE = 5 for cards with border-radius
    _add_rect(slide, MARGIN_H, Inches(1.9), col_w, Inches(5.0), surface, border, shape_type=5)
    y = Inches(2.1)
    for b in left_bullets:
        _add_text_box(slide, f"▸ {b}", MARGIN_H + Inches(0.15), y, col_w - Inches(0.3), Inches(0.5),
                      font_size=13, color=text_sec)
        y += Inches(0.48)

    rx = MARGIN_H + col_w + Inches(0.3)
    _add_rect(slide, rx, Inches(1.9), col_w, Inches(5.0), surface, border, shape_type=5)
    y = Inches(2.1)
    for b in right_bullets:
        _add_text_box(slide, f"▸ {b}", rx + Inches(0.15), y, col_w - Inches(0.3), Inches(0.5),
                      font_size=13, color=text_sec)
        y += Inches(0.48)
    _add_notes(slide, s.notes)


def _render_data(slide, s: SlideContent, colors: dict) -> None:
    bg = colors.get("gradient_bg", colors.get("background", "#1a1a2e"))
    primary = colors.get("primary", "#7c3aed")
    accent = colors.get("accent", "#f59e0b")
    text_sec = colors.get("text_secondary", "#94a3b8")
    surface = colors.get("surface", colors.get("card_bg", "#1e293b"))
    border = colors.get("border", "#475569")

    _apply_gradient_bg(slide, bg)
    _add_text_box(slide, s.title, MARGIN_H, MARGIN_V, CONTENT_W, Inches(0.9), 27, True, primary)
    # 80px × 3px accent line
    _add_rect(slide, MARGIN_H, Inches(1.6), Inches(0.83), Inches(0.03), accent)

    col_w = (CONTENT_W - Inches(0.3)) / 2
    # ROUNDED_RECTANGLE with 2px border (HTML: border-radius:16px, border:2px)
    _add_rect(slide, MARGIN_H, Inches(1.9), col_w, Inches(4.5), surface, border, shape_type=5, line_width=2)
    _add_text_box(slide, s.key_stat, MARGIN_H, Inches(2.5), col_w, Inches(1.5),
                  font_size=54, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _add_text_box(slide, s.key_stat_label, MARGIN_H, Inches(4.1), col_w, Inches(0.6),
                  font_size=12, color=text_sec, align=PP_ALIGN.CENTER)

    rx = MARGIN_H + col_w + Inches(0.3)
    y = Inches(2.0)
    for b in s.bullets:
        _add_text_box(slide, f"▸ {b}", rx, y, col_w, Inches(0.5), font_size=12, color=text_sec)
        y += Inches(0.5)
    _add_notes(slide, s.notes)


def _render_closing(slide, s: SlideContent, colors: dict) -> None:
    bg = colors.get("gradient_bg", colors.get("background", "#1a1a2e"))
    primary = colors.get("primary", "#7c3aed")
    accent = colors.get("accent", "#f59e0b")
    text_sec = colors.get("text_secondary", "#94a3b8")

    _apply_gradient_bg(slide, bg)
    # CONCLUSION label (18px → 14pt, HTML: letter-spacing:4px simulated with spaces)
    _add_text_box(slide, "C O N C L U S I O N", MARGIN_H, Inches(1.8), CONTENT_W, Inches(0.5),
                  font_size=14, bold=False, color=accent, align=PP_ALIGN.CENTER)
    _add_text_box(slide, s.title, MARGIN_H, Inches(2.4), CONTENT_W, Inches(1.6),
                  font_size=36, bold=True, color=primary, align=PP_ALIGN.CENTER)
    if s.subtitle:
        _add_text_box(slide, s.subtitle, MARGIN_H, Inches(4.0), CONTENT_W, Inches(1.2),
                      font_size=15, color=text_sec, align=PP_ALIGN.CENTER)
    # 60px × 4px accent line
    _add_rect(slide, SLIDE_W / 2 - Inches(0.31), Inches(5.5), Inches(0.63), Inches(0.04), accent)
    _add_notes(slide, s.notes)


_RENDERERS = {
    "cover": _render_cover,
    "section": _render_section,
    "content": _render_content,
    "two_column": _render_two_column,
    "data": _render_data,
    "closing": _render_closing,
}


def export_pptx(outline: PresentationOutline, colors: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    for slide_content in outline.slides:
        slide = prs.slides.add_slide(blank_layout)
        renderer = _RENDERERS.get(slide_content.slide_type, _render_content)
        renderer(slide, slide_content, colors)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out.absolute())
