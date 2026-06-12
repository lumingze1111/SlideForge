"""
pptx_converter.py — SlideForge HTML → PPTX 转换器

使用基于 Playwright 的测量 + python-pptx OOXML 装配流水线，
保证生成的 PPTX 与 HTML 预览高度一致。

用法：
    from slideforge.pptx_converter import convert_html_to_pptx
    convert_html_to_pptx("slides.html", "output.pptx")
"""

import json
import re
import tempfile
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def convert_html_to_pptx(
    html_path: str,
    output_path: str,
    embed_fonts: bool = False,
    verbose: bool = True,
    validate_gradients: bool = True,
    max_validate_attempts: int = 3,
    screenshot_mode: bool = False,
) -> str:
    """将 HTML 幻灯片文件转换为 PPTX。

    Args:
        html_path: HTML 文件路径（含 data-pptx-slide 属性的 .slide 元素）。
        output_path: 输出 .pptx 文件路径。
        embed_fonts: 是否嵌入字体（初始版本暂不支持）。
        verbose: 是否打印详细日志。
        validate_gradients: 装配后校验渐变是否与 HTML 匹配，不匹配则直接修补 PPTX。
        max_validate_attempts: 渐变校验循环最大尝试次数。
        screenshot_mode: 截图模式 —— 整页截图作背景，仅渲染透明文字叠层。

    Returns:
        输出 .pptx 的绝对路径。
    """
    html_path = Path(html_path).resolve()
    out_path = Path(output_path).resolve()

    if not html_path.exists():
        raise FileNotFoundError(f"HTML 文件不存在: {html_path}")

    # ── 提取演讲者备注 ───────────────────────────────────────────────
    notes_map = _extract_notes(html_path)

    # ── 预处理：确保每个 .slide 都带 data-pptx-slide ─────────────────
    measure_html = html_path
    cleanup_html: Path | None = None
    try:
        from tools.validate_gradients import ensure_pptx_slide_markers
        fixed, added = ensure_pptx_slide_markers(html_path)
        if fixed != html_path:
            measure_html = fixed
            cleanup_html = fixed
            if verbose and added:
                print(f"[converter] HTML 预处理: 已为 {added} 个 .slide 标注 data-pptx-slide")
    except ImportError:
        pass

    # ── Stage 1-2: 测量 → 装配 ─────────────────────────────────────
    from slideforge.pptx_engine.measure import measure
    from slideforge.pptx_engine.assemble import assemble

    if verbose:
        print("[converter] 正在测量 HTML 结构...")
    t0 = time.perf_counter()

    try:
        with tempfile.TemporaryDirectory(prefix="sf_pptx_") as tmp:
            tmp_dir = Path(tmp)
            meas_json = tmp_dir / "measurements.json"
            intermediate = tmp_dir / "no_fonts.pptx"

            # 测量
            measurement = measure(measure_html, meas_json,
                                  no_screenshots=not screenshot_mode,
                                  screenshot_mode=screenshot_mode,
                                  verbose=verbose)
            t1 = time.perf_counter()
            if verbose:
                print(f"[converter] 测量完成: {len(measurement.get('slides', []))} 页, "
                      f"{t1 - t0:.2f}s")

            # 装配
            if verbose:
                print("[converter] 正在装配 OOXML...")
            assemble(measurement, intermediate, screenshot_mode=screenshot_mode)
            t2 = time.perf_counter()
            if verbose:
                print(f"[converter] 装配完成: {t2 - t1:.2f}s")

            # 复制到输出路径
            import shutil
            shutil.copyfile(intermediate, out_path)

        # ── 注入演讲者备注 ──────────────────────────────────────────
        if notes_map:
            _inject_notes(str(out_path), notes_map, verbose=verbose)

        # ── 整体格式校验（截图模式下跳过渐变校验，截图已捕获渐变）─────
        if validate_gradients and not screenshot_mode:
            _run_format_validation(measure_html, out_path, measurement,
                                   max_validate_attempts, verbose)
    finally:
        if cleanup_html is not None and cleanup_html.exists():
            try:
                cleanup_html.unlink()
            except OSError:
                pass

    if verbose:
        print(f"[converter] ✓ PPTX 已生成: {out_path} "
              f"({out_path.stat().st_size:,} B)")

    return str(out_path.absolute())


def _extract_notes(html_path: Path) -> dict[int, str]:
    """从 HTML 中提取每张 slide 的演讲者备注。

    查找每个 [data-pptx-slide] 元素上的 data-notes 属性。
    返回 {1-based_index: notes_text}。
    """
    html = html_path.read_text(encoding="utf-8")
    pattern = re.compile(
        r'<div[^>]*?data-pptx-slide[^>]*?data-notes\s*=\s*"([^"]*)"',
        re.IGNORECASE,
    )
    notes_map = {}
    for i, m in enumerate(pattern.finditer(html), start=1):
        notes = m.group(1)
        # HTML 实体解码
        notes = notes.replace("&#39;", "'").replace("&quot;", '"').replace("&amp;", "&")
        notes = notes.replace("&lt;", "<").replace("&gt;", ">")
        notes = notes.replace("\\n", "\n")
        if notes.strip():
            notes_map[i] = notes
    return notes_map


def _inject_notes(pptx_path: str, notes_map: dict[int, str], verbose: bool = False) -> None:
    """向 PPTX 文件中注入演讲者备注。"""
    prs = Presentation(pptx_path)
    injected = 0
    for idx, slide in enumerate(prs.slides, start=1):
        notes_text = notes_map.get(idx)
        if not notes_text:
            continue
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = notes_text
            injected += 1
        except Exception:
            pass
    if injected:
        prs.save(pptx_path)
        if verbose:
            print(f"[converter] 已注入 {injected} 页演讲者备注")


def _run_format_validation(
    html_path: Path, out_path: Path, measurement: dict,
    max_attempts: int, verbose: bool,
) -> None:
    """对生成的 PPTX 做完整格式校验：渐变背景 + 文本样式 + 形状元素必须与 HTML 匹配。
    不匹配的项直接补丁 PPTX 的 slide XML，循环校验直至全部一致或超过尝试次数。
    """
    try:
        from tools.validate_gradients import (
            extract_html_gradients, extract_pptx_gradients,
            gradients_match, patch_pptx_gradients,
        )
        from tools.validate_format import (
            diff_decks, is_clean, collect_style_fixes, patch_pptx_styles,
        )
    except ImportError:
        if verbose:
            print("[converter] 校验模块未找到，跳过")
        return

    for attempt in range(1, max_attempts + 1):
        # 1. 渐变补丁
        expected_grads = extract_html_gradients(html_path)
        actual_grads = extract_pptx_gradients(out_path)
        if len(actual_grads) > len(expected_grads):
            actual_grads = actual_grads[:len(expected_grads)]
        elif len(actual_grads) < len(expected_grads):
            actual_grads += [None] * (len(expected_grads) - len(actual_grads))
        grad_fixes = {
            i: expected_grads[i]
            for i, (h, p) in enumerate(zip(expected_grads, actual_grads))
            if expected_grads[i] is not None and not gradients_match(h, p)
        }
        if grad_fixes:
            if verbose:
                print(f"[converter] 第 {attempt} 轮：修补渐变 "
                      f"{sorted(i + 1 for i in grad_fixes)}")
            patch_pptx_gradients(out_path, grad_fixes)

        # 2. 整体格式 diff
        report = diff_decks(measurement, out_path)
        t = report['totals']

        if is_clean(report):
            if verbose:
                print(f"[converter] ✓ 格式校验通过：{t['matched']}/{t['records']} "
                      f"records, {t['style_issues']} 样式, {t.get('pos_issues', 0)} 位置, "
                      f"{t['spurious']} 多余")
            return

        if verbose:
            issue_lines = [
                f"slide {s['index']}: {len(s['unmatched'])} 缺失, "
                f"{len(s['style_issues'])} 样式, {len(s.get('pos_issues', []))} 位置, "
                f"{len(s['spurious'])} 多余"
                for s in report['slides']
                if s['unmatched'] or s['style_issues'] or s.get('pos_issues') or s['spurious']
            ]
            print(f"[converter] 第 {attempt}/{max_attempts} 轮校验问题：")
            for line in issue_lines[:5]:
                print(f"    {line}")
            if len(issue_lines) > 5:
                print(f"    ... 共 {len(issue_lines)} 张幻灯片有问题")

        # 3. 样式补丁
        style_fixes = collect_style_fixes(report, measurement, out_path)
        if style_fixes:
            n = patch_pptx_styles(out_path, style_fixes)
            if verbose:
                print(f"[converter]   补丁了 {n} 个 run")
            continue

        # 没有可补丁的项 → 退出
        if not grad_fixes:
            break

    # 最终汇报
    if verbose:
        report = diff_decks(measurement, out_path)
        t = report['totals']
        if is_clean(report):
            print(f"[converter] ✓ 格式校验通过：{t['matched']}/{t['records']}")
        else:
            print(f"[converter] ⚠ 格式校验残留：{t['unmatched']} 缺失, "
                  f"{t['style_issues']} 样式, {t.get('pos_issues', 0)} 位置, "
                  f"{t['spurious']} 多余")
