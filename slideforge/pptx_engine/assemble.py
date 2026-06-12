"""assemble.py — 把 measurement JSON 装配成 pptx。

Usage:
    python assemble.py <measurement.json> <out.pptx>

设计：
- 标准 16:9 幻灯片 = 13.333" × 7.5" = 12192000 × 6858000 EMU
- 测量视口 1920×1080 → 1 CSS px = 6350 EMU, 元素/字体 1.5× 以中心缩放, 位置保持比例
- pptx 内部直接走低层 lxml 操作 spPr / txBody，避开 python-pptx 高层 API 的限制
"""
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W_EMU = 12192000      # 13.333"
SLIDE_H_EMU = 6858000       # 7.5"
PX_TO_EMU = SLIDE_W_EMU / SLIDE_W_PX  # 6350
PX_TO_PT = 0.5 * 1.5                    # 0.75 — 字体1.5倍(1920→13.333"空间基值0.5×1.5)
SIZE_SCALE = 1.5                         # 元素尺寸以中心为基准缩放

from slideforge.pptx_engine.embed_fonts import (
    family_alias_map, weighted_family_map, cjk_typefaces, cjk_for_style, style_of_typeface,
)
from slideforge.pptx_engine.text_utils import is_cjk_text

# 字体映射全部从 embed_fonts.FONT_PLAN 派生：
# - FONT_FALLBACKS: CSS 名 → OOXML typeface
# - CJK_FONTS:      标记为 CJK 的 typeface 集合
# FONT_PLAN 是运行时填充的（font_resolver 按需解析），所以 convert.py 在 resolve
# 之后必须调 refresh_font_plan_caches() 把这三个 module-level 缓存重新派生一次。
FONT_FALLBACKS: dict[str, str] = {}
WEIGHTED_FONT_FALLBACKS: dict[tuple[str, int, bool], str] = {}
CJK_FONTS: set[str] = set()
_CJK_ALIAS_SET: set[str] = set()


def refresh_font_plan_caches():
    """font_resolver 改完 FONT_PLAN 之后调一次，让 first_font / cjk_font 看到新条目。"""
    global FONT_FALLBACKS, WEIGHTED_FONT_FALLBACKS, CJK_FONTS, _CJK_ALIAS_SET
    FONT_FALLBACKS = family_alias_map()
    WEIGHTED_FONT_FALLBACKS = weighted_family_map()
    CJK_FONTS = cjk_typefaces()
    _CJK_ALIAS_SET = {name.lower() for name, tf in FONT_FALLBACKS.items() if tf in CJK_FONTS}


# import 时跑一次（FONT_PLAN 可能此时已被预填，例如 embed CLI 单独跑）
refresh_font_plan_caches()

def parse_text_shadow(value: str):
    """解析 CSS text-shadow，返回 (dx_px, dy_px, blur_px, (r,g,b,a)) 或 None。
    多层 shadow 取第一层。值如：
        "rgba(229, 57, 42, 1) 5px 5px 0px"
        "5px 5px rgb(0, 0, 0)"
        "rgb(244, 208, 63) 4px 4px 0px, rgb(15, 27, 61) 8px 8px 0px"  ← rgb() 内有逗号
    """
    if not value or value == "none":
        return None
    # 按顶层逗号切分（跳过括号内逗号）：computed value 的颜色经常是 rgb(R, G, B)，
    # 直接 split(",", 1) 会把 "rgb(244" 拆出来当第一层
    first = ""
    depth = 0
    for ch in value:
        if ch == "(":
            depth += 1
            first += ch
        elif ch == ")":
            depth -= 1
            first += ch
        elif ch == "," and depth == 0:
            break
        else:
            first += ch
    # 把可能在前后的 rgb()/rgba() 抠出来
    rgba_m = re.search(r"rgba?\(([^)]+)\)", first)
    color_rgba = (0, 0, 0, 1.0)
    if rgba_m:
        parts = [p.strip() for p in rgba_m.group(1).split(",")]
        if len(parts) >= 3:
            color_rgba = (int(float(parts[0])), int(float(parts[1])), int(float(parts[2])),
                          float(parts[3]) if len(parts) >= 4 else 1.0)
        first = re.sub(r"rgba?\([^)]+\)", "", first)
    nums = [float(m.group(1)) for m in re.finditer(r"(-?\d+\.?\d*)px", first)]
    if len(nums) < 2:
        return None
    dx, dy = nums[0], nums[1]
    blur = nums[2] if len(nums) >= 3 else 0.0
    return (dx, dy, blur, color_rgba)


def parse_rgb(s: str):
    """返回 (r,g,b) 元组；丢弃 alpha。要 alpha 请用 parse_rgba。"""
    return parse_rgba(s)[:3]


def _clamp_byte(value: float) -> int:
    return max(0, min(255, int(round(value))))


def _parse_css_alpha(value: str | None) -> float:
    if value is None or value == "":
        return 1.0
    v = str(value).strip()
    try:
        if v.endswith("%"):
            return max(0.0, min(1.0, float(v[:-1]) / 100.0))
        return max(0.0, min(1.0, float(v)))
    except ValueError:
        return 1.0


def _parse_css_rgb_component(value: str, srgb_unit: bool = False) -> int:
    v = str(value).strip()
    if v.lower() == "none":
        return 0
    try:
        if v.endswith("%"):
            return _clamp_byte(float(v[:-1]) * 2.55)
        n = float(v)
    except ValueError:
        return 0
    if srgb_unit:
        return _clamp_byte(n * 255.0)
    return _clamp_byte(n)


def parse_rgba(s: str):
    """返回 (r,g,b,a) 其中 a 是 0.0–1.0 浮点；缺省 1.0。"""
    if not s:
        return (0, 0, 0, 1.0)
    value = str(s).strip()
    if value in ("transparent", "rgba(0, 0, 0, 0)"):
        return (0, 0, 0, 0.0)
    if value.startswith("#"):
        hex_v = value[1:]
        if len(hex_v) in (3, 4):
            hex_v = "".join(ch * 2 for ch in hex_v)
        if len(hex_v) in (6, 8):
            try:
                r = int(hex_v[0:2], 16)
                g = int(hex_v[2:4], 16)
                b = int(hex_v[4:6], 16)
                a = int(hex_v[6:8], 16) / 255.0 if len(hex_v) == 8 else 1.0
                return (r, g, b, a)
            except ValueError:
                return (0, 0, 0, 1.0)

    m = re.match(r"rgba?\(([^)]+)\)", value)
    if m:
        body = m.group(1).strip()
        if "," in body:
            parts = [p.strip() for p in body.split(",")]
            rgb_parts = parts[:3]
            alpha_part = parts[3] if len(parts) >= 4 else None
        else:
            left, sep, right = body.partition("/")
            rgb_parts = [p for p in left.split() if p]
            alpha_part = right.strip() if sep else None
        if len(rgb_parts) >= 3:
            return (
                _parse_css_rgb_component(rgb_parts[0]),
                _parse_css_rgb_component(rgb_parts[1]),
                _parse_css_rgb_component(rgb_parts[2]),
                _parse_css_alpha(alpha_part),
            )

    m = re.match(r"color\(\s*srgb\s+([^)]+)\)", value)
    if m:
        body = m.group(1).strip()
        left, sep, right = body.partition("/")
        parts = [p for p in left.split() if p]
        if len(parts) >= 3:
            return (
                _parse_css_rgb_component(parts[0], srgb_unit=True),
                _parse_css_rgb_component(parts[1], srgb_unit=True),
                _parse_css_rgb_component(parts[2], srgb_unit=True),
                _parse_css_alpha(right.strip() if sep else None),
            )

    return (0, 0, 0, 1.0)


GENERIC_FONT_KEYWORDS = {
    # CSS 通用字体族关键字，不算"具体字体"，遇到要跳过
    "serif", "sans-serif", "monospace", "cursive", "fantasy",
    "system-ui", "ui-serif", "ui-sans-serif", "ui-monospace",
    "math", "emoji", "fangsong",
    # 平台系统 UI 字体别名：CSS 端是关键字（解析为当前 OS 系统字体），
    # 但 OOXML typeface= 写它们 PowerPoint/WPS 找不到真实字体，会落到默认 fallback
    # 并对 b="1" run 做 faux-bold 双绘（字符堆积）。当作 generic 跳过即可。
    "-apple-system", "blinkmacsystemfont", "-webkit-system-font",
}

# 最后兜底字体（CSS 没声明、或全是 generic 关键字时用）。Calibri 是 PowerPoint 自 2007
# 起的默认西文字体，Windows + Office 跨平台都自带；其它 viewer 上若缺则各自系统再 fallback。
DEFAULT_LATIN_FALLBACK = "Calibri"


def first_font(font_family: str) -> str:
    """从 CSS font-family 字符串里挑第一项（latin 用），去引号。
    优先返回 FONT_FALLBACKS 里映射到的 OOXML typeface；跳过 generic 关键字。"""
    items = [x.strip().strip('"').strip("'") for x in font_family.split(",")]
    for it in items:
        if not it or it.lower() in GENERIC_FONT_KEYWORDS:
            continue
        # 大小写不敏感查 FONT_FALLBACKS（处理 CSS 大小写差异）
        if it in FONT_FALLBACKS:
            return FONT_FALLBACKS[it]
        if it.lower() in FONT_FALLBACKS:
            return FONT_FALLBACKS[it.lower()]
        return it  # 用户用了我们没装的字体，原名透传（运行时回退到系统）
    return items[0] if items else DEFAULT_LATIN_FALLBACK


def _normalize_weight_value(weight) -> int:
    try:
        return int(float(weight))
    except (TypeError, ValueError):
        if str(weight).lower() in ("bold", "bolder"):
            return 700
        return 400


def _is_italic_value(style) -> bool:
    return str(style or "").lower() == "italic"


def first_font_for_run(font_family: str, font_weight, font_style) -> tuple[str, bool]:
    """Return (OOXML typeface, exact_weight_face).

    When font_resolver embedded exact source weights, use the requested run
    weight to select the matching typeface dynamically.
    """
    items = [x.strip().strip('"').strip("'") for x in font_family.split(",")]
    weight = _normalize_weight_value(font_weight)
    italic = _is_italic_value(font_style)

    for it in items:
        if not it or it.lower() in GENERIC_FONT_KEYWORDS:
            continue
        key = (it.lower(), weight, italic)
        if key in WEIGHTED_FONT_FALLBACKS:
            return WEIGHTED_FONT_FALLBACKS[key], True

        candidates = [
            (abs(w - weight), w > weight, w, typeface)
            for (fam, w, it_italic), typeface in WEIGHTED_FONT_FALLBACKS.items()
            if fam == it.lower() and it_italic == italic
        ]
        if candidates:
            return min(candidates)[3], True

    return first_font(font_family), False


def cjk_font(font_family: str, latin_name: str) -> str:
    """返回该 run 应使用的 East Asian 字体。

    决策顺序（全部从 FONT_PLAN 派生，无硬编码）：
    1. CSS family 列表里显式列了 CJK 字体（或其 alias） → 用映射后的 typeface
    2. 否则按 latin 字体的 style 配对：latin serif → CJK serif，latin sans/mono → CJK sans
    3. style 未知时默认 sans CJK
    """
    items = [x.strip().strip('"').strip("'") for x in font_family.split(",")]
    for it in items:
        if it.lower() in _CJK_ALIAS_SET:
            return FONT_FALLBACKS.get(it, FONT_FALLBACKS.get(it.lower(), it))
    # 没有显式 CJK：按 latin 字体的 style 选配对的 CJK
    latin_style = style_of_typeface(latin_name)
    return cjk_for_style(latin_style)


def px_to_emu(px: float, scale: float = 1.0) -> int:
    """CSS px → EMU。scale=1.0 用于位置 (x,y)，scale=SIZE_SCALE 用于宽高/边距/线宽。"""
    return int(round(px * PX_TO_EMU * scale))


def _scaled_rect(rx: float, ry: float, rw: float, rh: float) -> tuple[int, int, int, int]:
    """返回 (x_emu, y_emu, w_emu, h_emu)，元素以中心为基准缩放 SIZE_SCALE 倍。

    尺寸按 SIZE_SCALE 放大，同时位置向左上偏移，使元素视觉中心保持不变，
    避免元素向右下偏移。
    """
    offset = (SIZE_SCALE - 1.0) / 2.0  # 0.25
    x = px_to_emu(rx - rw * offset)
    y = px_to_emu(ry - rh * offset)
    w = px_to_emu(rw, SIZE_SCALE)
    h = px_to_emu(rh, SIZE_SCALE)
    return _clamp_to_slide(x, y, w, h)


def _clamp_to_slide(x: int, y: int, w: int, h: int) -> tuple[int, int, int, int]:
    """确保元素不超出幻灯片边界。优先移动位置保留完整尺寸，仅当元素宽/高超过
    幻灯片时裁剪尺寸。"""
    sw, sh = SLIDE_W_EMU, SLIDE_H_EMU

    # 宽度超过幻灯片 → 裁剪
    if w > sw:
        w = sw

    # 水平：优先移动位置使元素完整可见
    if x < 0:
        x = 0
    elif x + w > sw:
        x = sw - w

    # 高度超过幻灯片 → 裁剪
    if h > sh:
        h = sh

    # 垂直：优先移动位置使元素完整可见
    if y < 0:
        y = 0
    elif y + h > sh:
        y = sh - h

    return x, y, max(0, w), max(0, h)


def _resolve_rect(rec, fallback_w=None, fallback_h=None):
    """返回 (x_emu, y_emu, w_emu, h_emu)，优先使用 _adjusted_rect。

    _adjusted_rect = (x, y, w, h) CSS px 绝对坐标，由 Layout Agent 输出。
    位置转 EMU 用 scale=1.0，尺寸转 EMU 用 SIZE_SCALE（1.5×）。

    fallback_w/fallback_h 覆盖 _scaled_rect 的尺寸参数
    （用于 add_text_box 已单独算出 textbox 尺寸的情况）。
    """
    adjusted = rec.get("_adjusted_rect")
    if adjusted:
        x = px_to_emu(adjusted[0])
        y = px_to_emu(adjusted[1])
        w = px_to_emu(adjusted[2], SIZE_SCALE)
        h = px_to_emu(adjusted[3], SIZE_SCALE)
        return _clamp_to_slide(x, y, w, h)
    r = rec["rect"]
    return _scaled_rect(r["x"], r["y"],
                        fallback_w if fallback_w is not None else r["w"],
                        fallback_h if fallback_h is not None else r["h"])


def _find_blank_layout(prs):
    """从 prs.slide_layouts 找空白布局，避免硬编码下标 6。

    顺序：
    1. layout.name == 'Blank'（python-pptx 默认模板里就叫这个名字）
    2. 没有 placeholder 的 layout（结构上等同空白）
    3. 兜底用最后一个 layout
    """
    layouts = list(prs.slide_layouts)
    for layout in layouts:
        if (layout.name or "").strip().lower() == "blank":
            return layout
    for layout in layouts:
        if len(list(layout.placeholders)) == 0:
            return layout
    return layouts[-1]


def _text_max_font_size(rec) -> float:
    runs = rec.get("runs", []) or []
    return max((float(run.get("fontSize", 16) or 16) for run in runs), default=16)


def _has_explicit_break(rec) -> bool:
    """rec 的 runs 里是否有显式 `<br>` / `\n` 分行。"""
    for run in (rec.get("runs", []) or []):
        text = run.get("text", "") or ""
        if run.get("linebreak") or (text.strip() and "\n" in text):
            return True
    return False


def _text_is_single_line(rec, max_fs: float) -> bool:
    """判断 record 是不是单行文本。
    用 BCR.h 减掉 CSS padding（垂直方向）得到内容高度，再跟 max_fs*1.8 比。
    不扣 padding 的话，带 padding 的短标签（button / badge / pill）会被误判成多行，
    走 wrap=square 让 PPT 度量稍宽时把短词切开（如 "SELECT" → "SELEC"/"T"）。
    """
    r = rec["rect"]
    style = rec.get("style", {}) or {}
    pad_v = (style.get("paddingTop", 0) or 0) + (style.get("paddingBottom", 0) or 0)
    content_h = max(0, r["h"] - pad_v)
    return content_h < max_fs * 1.8 and not _has_explicit_break(rec)


def _text_box_size_px(rec, max_fs: float, is_single_line: bool) -> tuple[float, float]:
    """textbox 几何 = HTML BCR × 1.12 余量。

    PPT 与浏览器的字体度量存在系统性差异（同字号下 PPT 字宽略大），
    在 1.5× 整体缩放后差异被放大导致文字溢出。给 12% 额外空间补偿，
    远小于旧版 1.3×–1.4× 的垂直膨胀，不会造成段落间叠压。
    """
    r = rec["rect"]
    return r["w"] * 1.12, r["h"] * 1.12


def _prepare_text_layouts(records):
    """缓存每个 text record 的 layout 计算（max_fs / is_single_line / textbox 尺寸）。
    一次计算两次复用：text_box_overlap_warnings 在 self_check 也会读 _pptx_text_layout。"""
    text_records = [rec for rec in records if rec.get("kind") == "text" and rec.get("rect")]
    for rec in text_records:
        max_fs = _text_max_font_size(rec)
        is_single_line = _text_is_single_line(rec, max_fs)
        w_px, h_px = _text_box_size_px(rec, max_fs, is_single_line)
        rec["_pptx_text_layout"] = {
            "max_fs": max_fs,
            "is_single_line": is_single_line,
            "w_px": w_px,
            "h_px": h_px,
        }


def make_rgb(rgb):
    r, g, b = rgb
    return RGBColor(r, g, b)


def _apply_rotation(shape, rec):
    """如果 record 上有非零 rotation，重写 shape 的 xfrm：
    - off/ext 用 naturalSize（未旋转的元素本尺寸）+ AABB 中心点（rect 是浏览器返回的旋转后 AABB）
    - 加 rot 属性（1/60000 度）
    会让 PPT 按"先放在 AABB 中心、再旋转"的方式还原 HTML 的视觉效果。
    """
    rot_deg = rec.get("rotation") or 0.0
    if abs(rot_deg) < 0.5:
        return
    nat = rec.get("naturalSize") or {}
    nat_w = float(nat.get("w") or 0)
    nat_h = float(nat.get("h") or 0)
    if nat_w <= 0 or nat_h <= 0:
        return
    rect = rec["rect"]
    cx_px = rect["x"] + rect["w"] / 2.0
    cy_px = rect["y"] + rect["h"] / 2.0
    new_x = px_to_emu(cx_px) - px_to_emu(nat_w, SIZE_SCALE) // 2
    new_y = px_to_emu(cy_px) - px_to_emu(nat_h, SIZE_SCALE) // 2
    new_w = px_to_emu(nat_w, SIZE_SCALE)
    new_h = px_to_emu(nat_h, SIZE_SCALE)
    new_x, new_y, new_w, new_h = _clamp_to_slide(new_x, new_y, new_w, new_h)
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return
    # OOXML rot 单位 = 1/60000 度，正数 = 顺时针
    rot_units = int(round(rot_deg * 60000)) % (360 * 60000)
    xfrm.set("rot", str(rot_units))
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    if off is not None:
        off.set("x", str(new_x))
        off.set("y", str(new_y))
    if ext is not None:
        ext.set("cx", str(new_w))
        ext.set("cy", str(new_h))


def _parse_gradient_from_css(bg_image: str) -> dict | None:
    """从 CSS backgroundImage 计算值中解析渐变。

    输入: 'linear-gradient(135deg, rgb(15, 23, 42), rgb(30, 41, 59))'
    返回: {'angle': 135, 'stops': [((15,23,42), 0), ((30,41,59), 1000)]}

    支持 linear-gradient，单位 deg/rad/turn/grad。
    """
    if not bg_image or bg_image == "none":
        return None
    m = re.match(
        r"(?:-webkit-)?(linear-gradient)\s*\(", bg_image.strip(), re.IGNORECASE
    )
    if not m:
        return None

    # 提取括号体（处理嵌套括号）
    body = _extract_paren_body(bg_image.strip(), m.end() - 1)
    if body is None:
        return None

    tokens = _split_gradient_body(body)
    if not tokens:
        return None

    angle = 0
    stop_offset = 0

    # 检查第一个 token 是不是角度
    first = tokens[0].strip()
    angle_m = re.match(
        r"([+-]?\d+(?:\.\d+)?)\s*(deg|rad|turn|grad)\s*$", first, re.IGNORECASE
    )
    if angle_m:
        val = float(angle_m.group(1))
        unit = angle_m.group(2).lower()
        if unit == "rad":
            val = val * 180 / 3.1415926535
        elif unit == "turn":
            val = val * 360
        elif unit == "grad":
            val = val * 0.9
        angle = val % 360
        stop_offset = 1

    # 解析 color stops
    stops = []
    raw_stops = tokens[stop_offset:]
    explicit_positions = []

    for raw in raw_stops:
        s = raw.strip()
        # Matches: rgb(r,g,b) position%, rgba(r,g,b,a) position%, hex position%
        color_stop_m = re.match(
            r"^(rgba?\([^)]+\)|#[0-9a-fA-F]+)\s*(.*)", s, re.IGNORECASE
        )
        if not color_stop_m:
            continue
        color_str = color_stop_m.group(1)
        pos_str = color_stop_m.group(2).strip()

        rgba = parse_rgba(color_str)
        pos = None
        if pos_str:
            pct_m = re.match(r"([+-]?\d+(?:\.\d+)?)\s*%", pos_str)
            if pct_m:
                pos = int(float(pct_m.group(1)) * 10)

        stops.append((rgba, pos, color_str))
        if pos is not None:
            explicit_positions.append(len(stops) - 1)

    if len(stops) < 2:
        return None

    # Fill in missing positions
    if explicit_positions:
        last_filled = None
        next_filled = None
        for i, (_, pos, _) in enumerate(stops):
            if pos is not None:
                if last_filled is None:
                    last_filled = i
                next_filled = None
                for j in range(i + 1, len(stops)):
                    if stops[j][1] is not None:
                        next_filled = j
                        break
            else:
                if last_filled is not None and next_filled is not None:
                    l_pos = stops[last_filled][1]
                    r_pos = stops[next_filled][1]
                    frac = (i - last_filled) / (next_filled - last_filled)
                    stops[i] = (stops[i][0], int(l_pos + frac * (r_pos - l_pos)), stops[i][2])
                elif last_filled is not None:
                    stops[i] = (stops[i][0], stops[last_filled][1], stops[i][2])
                elif next_filled is not None:
                    stops[i] = (stops[i][0], stops[next_filled][1], stops[i][2])
                else:
                    stops[i] = (stops[i][0], 0, stops[i][2])
    else:
        n = len(stops)
        for i in range(n):
            if i == 0:
                stops[i] = (stops[i][0], 0, stops[i][2])
            elif i == n - 1:
                stops[i] = (stops[i][0], 1000, stops[i][2])
            else:
                stops[i] = (stops[i][0], int(i * 1000 / (n - 1)), stops[i][2])

    return {"angle": angle, "stops": [(s[0][:3], s[1]) for s in stops]}


def _extract_paren_body(s: str, open_pos: int) -> str | None:
    """从 open_pos 的 '(' 开始，提取匹配的括号体内容（不含外层括号）。
    正确处理嵌套括号。
    """
    depth = 0
    started = False
    start = -1
    for i, ch in enumerate(s):
        if i < open_pos:
            continue
        if ch == "(":
            if not started:
                started = True
                start = i + 1
            depth += 1
        elif ch == ")":
            depth -= 1
            if started and depth == 0:
                return s[start:i]
    return None


def _split_gradient_body(body: str) -> list[str]:
    """按顶层逗号切分 gradient body（跳过括号内逗号）。"""
    parts = []
    depth = 0
    cur = ""
    for ch in body:
        if ch == "(":
            depth += 1
            cur += ch
        elif ch == ")":
            depth -= 1
            cur += ch
        elif ch == "," and depth == 0:
            parts.append(cur)
            cur = ""
        else:
            cur += ch
    if cur.strip():
        parts.append(cur)
    return parts


def add_background(slide, rgb, background_image: str | None = None):
    """整页底色。支持渐变（通过 CSS backgroundImage）和纯色。"""
    gradient = _parse_gradient_from_css(background_image) if background_image else None
    if gradient:
        _apply_gradient_rect(slide, gradient)
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = make_rgb(rgb)


def _apply_gradient_rect(slide, gradient: dict) -> None:
    """用全页矩形 + 渐变填充实现背景渐变。

    相比直接操作 slide background OOXML（p:bg/p:bgPr），形状渐变在各 PPT 渲染器
    （PowerPoint / WPS / Keynote）中兼容性更好。
    """
    from pptx.enum.shapes import MSO_SHAPE

    # slide-level 背景留空（避免在渐变 rect 之外有色差，或渲染器把 slide bg
    # 当作"已存在的内容"叠加显示）。渐变矩形铺满 SLIDE_W × SLIDE_H，足以覆盖。
    fill = slide.background.fill
    fill.background()

    # 叠一个铺满整页的矩形，打上原生渐变
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W_EMU, SLIDE_H_EMU
    )
    rect.line.fill.background()   # 无边框

    # 3. 把 python-pptx 默认的 solidFill 替换成 gradFill
    spPr = rect._element.find(qn("p:spPr"))
    if spPr is None:
        return
    for child in list(spPr):
        if child.tag.endswith("solidFill") or child.tag.endswith("noFill"):
            spPr.remove(child)

    angle = gradient["angle"]
    ooxml_angle = int((angle - 90) % 360 * 60000)

    grad_fill = etree.SubElement(spPr, qn("a:gradFill"))
    grad_fill.set("rotWithShape", "1")

    gs_lst = etree.SubElement(grad_fill, qn("a:gsLst"))
    for (r, g, b), pos in gradient["stops"]:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        # OOXML pos is ST_PositiveFixedPercentage (0..100000, 1000ths of a percent).
        # Internal `pos` uses 0..1000 (10ths of a percent), so multiply by 100.
        gs.set("pos", str(int(pos) * 100))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", "{:02X}{:02X}{:02X}".format(r, g, b))

    lin = etree.SubElement(grad_fill, qn("a:lin"))
    lin.set("ang", str(ooxml_angle))
    lin.set("scaled", "1")

    # 4. 清除 python-pptx 加上的 p:style（QuickStyle），它的 fillRef 会覆盖 gradFill
    style_el = rect._element.find(qn("p:style"))
    if style_el is not None:
        rect._element.remove(style_el)

    # 5. 清除 txBody（矩形不需要文字框）
    txBody = rect._element.find(qn("p:txBody"))
    if txBody is not None:
        rect._element.remove(txBody)


def add_shape_box(slide, rec):
    """带 background / border 的非文本装饰节点。

    border 处理：每条边单独判断。
    - 全 4 边都有：画带边的矩形
    - 单侧或不对称：用 connectorStraightLine 单独画每条出现的边线
    """
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    x, y, w, h = _resolve_rect(rec)
    deco = rec.get("deco", {})

    sides = {
        "top":    (deco.get("borderTop"),    deco.get("borderTopWidth", 0)),
        "bottom": (deco.get("borderBottom"), deco.get("borderBottomWidth", 0)),
        "left":   (deco.get("borderLeft"),   deco.get("borderLeftWidth", 0)),
        "right":  (deco.get("borderRight"),  deco.get("borderRightWidth", 0)),
    }
    active_sides = [k for k, (present, _) in sides.items() if present]

    # 圆角分流：CSS border-radius >= 50% 或 >= min(w,h)/2 → 椭圆 / pill / 矩形 三选一
    # 接近方形 → OVAL；宽高悬殊 → ROUNDED_RECTANGLE + adj=0.5 (pill)；都不是 → RECTANGLE
    kind = _round_kind(deco.get("borderRadius", ""), r["w"], r["h"])
    prst = {"oval": MSO_SHAPE.OVAL,
            "pill": MSO_SHAPE.ROUNDED_RECTANGLE,
            "rect": MSO_SHAPE.RECTANGLE}[kind]
    is_round = (kind != "rect")

    side_color_keys = {
        "top": "borderTopColor",
        "bottom": "borderBottomColor",
        "left": "borderLeftColor",
        "right": "borderRightColor",
    }

    def border_rgba_for(side):
        return parse_rgba(deco.get(side_color_keys[side]) or deco.get("borderColor", "rgb(127,127,127)"))

    border_rgba = border_rgba_for("top")
    same_border_colors = all(border_rgba_for(side) == border_rgba for side in active_sides)
    bw_top = deco.get("borderTopWidth", 0) or 0
    bw_bottom = deco.get("borderBottomWidth", 0) or 0
    bw_left = deco.get("borderLeftWidth", 0) or 0
    bw_right = deco.get("borderRightWidth", 0) or 0
    widest = max(bw_top, bw_bottom, bw_left, bw_right)
    collapsed_border_box = (
        r["h"] <= (bw_top + bw_bottom + 0.5)
        or r["w"] <= (bw_left + bw_right + 0.5)
    )
    equal_border_widths = (
        abs(bw_top - bw_bottom) < 0.01
        and abs(bw_top - bw_left) < 0.01
        and abs(bw_top - bw_right) < 0.01
    )

    # Simple four-sided boxes should stay one editable PPT shape. Previously a
    # filled rectangle plus four line objects created selectable "background"
    # blocks inside every card.
    if same_border_colors and len(active_sides) == 4 and (not deco.get("hasBg") or (kind == "rect" and equal_border_widths)):
        shape = slide.shapes.add_shape(prst, x, y, w, h)
        if kind == "pill":
            shape.adjustments[0] = 0.5

        if deco.get("hasBg"):
            r_, g_, b_, a_ = parse_rgba(deco["bg"])
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r_, g_, b_)
            if a_ < 1.0:
                _set_fill_alpha(shape, a_)
        elif collapsed_border_box:
            # CSS border is drawn inside the border box. If the box is thinner
            # than its opposing borders, the visual result is a solid strip.
            shape.fill.solid()
            shape.fill.fore_color.rgb = make_rgb(border_rgba[:3])
            if border_rgba[3] < 1.0:
                _set_fill_alpha(shape, border_rgba[3])
        else:
            shape.fill.background()

        if collapsed_border_box and not deco.get("hasBg"):
            shape.line.fill.background()
        else:
            shape.line.color.rgb = make_rgb(border_rgba[:3])
            shape.line.width = Emu(px_to_emu(widest or 1, SIZE_SCALE))
            _set_line_alpha(shape, border_rgba[3])
        _apply_rotation(shape, rec)
        return shape

    # 如果有填充色：画形状（不带 border，border 单独画线）
    fill_shape = None
    if deco.get("hasBg"):
        r_, g_, b_, a_ = parse_rgba(deco["bg"])
        shape = slide.shapes.add_shape(prst, x, y, w, h)
        if kind == "pill":
            shape.adjustments[0] = 0.5  # 圆角半径 = 短边 50%，端点完整半圆
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r_, g_, b_)
        # 把 CSS alpha 转成 OOXML 的 a:alpha（单位 1/1000 百分比，100000 = 100%）
        if a_ < 1.0:
            _set_fill_alpha(shape, a_)
        shape.line.fill.background()
        _apply_rotation(shape, rec)
        fill_shape = shape

    # oval / pill（仅有填充无 4 边边框）：完了直接 return，跳过下面"按需画线"
    if is_round and deco.get("hasBg"):
        return fill_shape

    # 否则按需画线（border-top / border-bottom 等单侧情形）
    for side in active_sides:
        side_rgba = border_rgba_for(side)
        rgb = side_rgba[:3]
        alpha = side_rgba[3]
        bw = sides[side][1] or 1
        if side == "top":
            _add_line(slide, x, y, x + w, y, rgb, bw, alpha)
        elif side == "bottom":
            _add_line(slide, x, y + h, x + w, y + h, rgb, bw, alpha)
        elif side == "left":
            _add_line(slide, x, y, x, y + h, rgb, bw, alpha)
        elif side == "right":
            _add_line(slide, x + w, y, x + w, y + h, rgb, bw, alpha)


def _round_kind(border_radius: str, w_px: float, h_px: float) -> str:
    """border-radius + 元素宽高 → 'oval' | 'pill' | 'rect'。

    - 半径不够大（< min(w,h)/2 * 0.9 或 < 50%）→ rect
    - 半径够大 + 宽高接近方形（0.67 ≤ w/h ≤ 1.5）→ oval
    - 半径够大 + 宽高悬殊 → pill（OOXML roundRect adj=0.5，两端完整半圆 + 中间矩形）

    pill 阈值 1.5：低于这个宽高比的"扁形" oval 视觉上还能接受；超过就明显是 CSS pill 意图。
    """
    if not border_radius or border_radius == "0px":
        return "rect"
    s = str(border_radius).strip()
    big = False
    if s.endswith("%"):
        try:
            big = float(s[:-1]) >= 50.0   # 去掉最后 1 个 "%" 字符
        except ValueError:
            return "rect"
    elif s.endswith("px"):
        try:
            # 至少要达到 短边/2 × 0.9 才算"贴边圆角"
            big = float(s[:-2]) >= min(w_px, h_px) / 2 * 0.9
        except ValueError:
            return "rect"
    else:
        return "rect"
    if not big:
        return "rect"
    if w_px <= 0 or h_px <= 0:
        return "oval"
    ratio = w_px / h_px
    if 0.67 <= ratio <= 1.5:
        return "oval"
    return "pill"


def _set_fill_alpha(shape, alpha: float):
    """在 shape 的 solidFill 上加 a:alpha 子元素。alpha ∈ [0,1]。"""
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # OOXML a:alpha val 单位为 千分之一（100000 = 100%）
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(int(round(alpha * 100000))))


def _set_line_alpha(shape_or_line, alpha: float):
    """在 spPr/a:ln/a:solidFill/a:srgbClr 上加 a:alpha 子元素。alpha ∈ [0,1]。
    alpha >= 1.0 时直接返回（OOXML 默认就是不透明）。"""
    if alpha >= 1.0:
        return
    spPr = shape_or_line._element.find(qn("p:spPr"))
    if spPr is None:
        return
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return
    solid = ln.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(int(round(alpha * 100000))))


def _add_line(slide, x1, y1, x2, y2, color_rgb, width_px, alpha: float = 1.0):
    """画一条直线。x1/y1/x2/y2 已是 EMU。可选 alpha 支持半透明。"""
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = make_rgb(color_rgb)
    line.line.width = Emu(px_to_emu(width_px, SIZE_SCALE))
    _set_line_alpha(line, alpha)
    return line


def _normalize_runs(runs):
    """把测到的 runs 折叠成 OOXML 段落生成器要消费的 (kind, ...) 项序列。

    复刻 CSS 空白折叠：每个 run 内 \\s+ → 单空格，去掉首尾空白，过滤纯空白 run；
    跨 run 的"边界空格"用单独的 SPACE 项保留，避免 run 拼接后丢空格。

    返回 list[(kind, ...)], 其中 kind ∈ {"BREAK", "SPACE", "RUN"}：
    - ("BREAK", run)        显式 <br> / \\n
    - ("SPACE", None)       run 之间的边界空格
    - ("RUN", run, body)    实际文本 run
    """
    cleaned = []
    for i, run in enumerate(runs):
        text = run.get("text", "")
        if not text:
            continue
        if run.get("linebreak"):
            cleaned.append(("BREAK", run))
            continue
        collapsed = re.sub(r"\s+", " ", text)
        leading = collapsed.startswith(" ")
        trailing = collapsed.endswith(" ")
        body = collapsed.strip()
        if not body:
            # 纯空白 run：若非首末位，当作 run 间空格信号
            if cleaned and i != len(runs) - 1:
                cleaned.append(("SPACE", None))
            continue
        if leading and cleaned and cleaned[-1][0] != "BREAK":
            cleaned.append(("SPACE", None))
        cleaned.append(("RUN", run, body))
        if trailing and i != len(runs) - 1:
            cleaned.append(("SPACE", None))

    # 去掉相邻的重复 SPACE
    deduped = []
    for item in cleaned:
        if item[0] == "SPACE" and deduped and deduped[-1][0] == "SPACE":
            continue
        deduped.append(item)
    return deduped


def _fix_empty_paragraph_sizes(tf, style_font_size_px):
    """把空段（<a:p> 无 <a:r>）的 endParaRPr.sz 设成正文字号。

    不显式设的话，OOXML 空段会用 PPT 默认 18pt × 行距撑出 ≈43px 空行高，
    把后续内容推出 textbox 砸到下方相邻段。
    """
    fs_px = float(style_font_size_px or 16) or 16
    end_sz = max(100, int(round(fs_px * PX_TO_PT * 100)))
    for para in tf.paragraphs:
        if para._p.find(qn("a:r")) is None:
            endR = para._p.find(qn("a:endParaRPr"))
            if endR is None:
                endR = etree.SubElement(para._p, qn("a:endParaRPr"))
            endR.set("sz", str(end_sz))


def _emit_line_break(paragraph, style_font_size_px):
    """Insert an OOXML soft line break inside the current paragraph.

    Browser-wrapped lines are still one CSS paragraph. Emitting a new PPT
    paragraph adds paragraph-model spacing in PowerPoint/WPS and makes the
    visual line height looser than HTML.
    """
    br = etree.SubElement(paragraph._p, qn("a:br"))
    rPr = etree.SubElement(br, qn("a:rPr"))
    fs_px = float(style_font_size_px or 16) or 16
    rPr.set("sz", str(max(100, int(round(fs_px * PX_TO_PT * 100)))))


def add_text_box(slide, rec):
    """文本节点 → pptx textbox（多 run 富文本）。"""
    r = rec["rect"]
    # 若 text leaf 同时带背景色 / 边框（如 .bar.a 既是柱子又装着文字），
    # 先按 rec 的 deco 画一个 shape 垫底，再画文字框 —— 否则背景丢失。
    deco = rec.get("deco", {})
    has_decoration = deco.get("hasBg") or deco.get("borderTop") or deco.get("borderBottom") \
                     or deco.get("borderLeft") or deco.get("borderRight")
    if has_decoration:
        # 用一个仅含 deco 的合成 record 调 add_shape_box
        synth = {"rect": r, "deco": deco, "kind": "shape", "tag": rec.get("tag", "div")}
        add_shape_box(slide, synth)

    # 判断单行 / 多行：浏览器测得高度 < 1.8 × 主字号 → 单行
    layout = rec.get("_pptx_text_layout")
    if layout is None:
        max_fs = _text_max_font_size(rec)
        is_single_line = _text_is_single_line(rec, max_fs)
        w_px, h_px = _text_box_size_px(rec, max_fs, is_single_line)
    else:
        max_fs = layout["max_fs"]
        is_single_line = layout["is_single_line"]
        w_px = layout["w_px"]
        h_px = layout["h_px"]

    if w_px <= 0:
        w_px = 50.0
    if h_px <= 0:
        h_px = 20.0
    x, y, w, h = _resolve_rect(rec, w_px, h_px)

    tb = slide.shapes.add_textbox(x, y, w, h)
    _apply_rotation(tb, rec)
    tf = tb.text_frame
    # 何时禁用自动 wrap：
    # (1) 真单行（h < max_fs*1.8）— 短标签 PPT 度量稍宽时不应换行
    # (2) 紧排版多行（h ≈ (br_count+1) × line-height）— 每段恰好一行，
    #     PPT 度量稍宽时不应再切，会破坏作者用 <br> 的排版意图（如 "Title<br>Subtitle"）
    # 松排版（h 远大于 BR 分段所需）：段落本身要 word-wrap，必须 wrap=square 让 PPT 按宽切，
    # 否则长段落会撑成单行串过整张幻灯片（典型："<p>...long</p><br><br><p>...long</p>" 双段排版）
    runs_raw = rec.get("runs") or []
    br_count = sum(1 for run in runs_raw if run.get("linebreak"))
    has_explicit_break = br_count > 0
    # 计算"有效行距"：leaf 自身 lineHeight（容器 computed 值）vs runs 内最大字号节点的
    # lineHeight，取较大者。叶子 .fadelist-items (16px / lh 14.72px) 但 span 是 144px / lh 144px
    # 这种 case，按 leaf 算行距会让 144px 文字被压在 14.72px 行距里叠压
    style_for_lh = rec.get("style", {})
    leaf_fs = style_for_lh.get("fontSize", 16)
    effective_lh_px = _line_height_px(style_for_lh.get("lineHeight"), leaf_fs)
    effective_fs = leaf_fs
    for run in runs_raw:
        if run.get("linebreak"):
            continue
        rfs = run.get("fontSize")
        rlh = run.get("lineHeight")
        if rfs and rlh:
            rlh_px = _line_height_px(rlh, rfs)
            if rlh_px > effective_lh_px:
                effective_lh_px = rlh_px
                effective_fs = rfs
    expected_lines = (r["h"] / effective_lh_px) if effective_lh_px > 0 else 1.0
    is_tight_layout = expected_lines <= br_count + 1 + 0.7
    no_auto_wrap = is_single_line or (has_explicit_break and is_tight_layout)
    tf.word_wrap = not no_auto_wrap
    # textbox 几何 = HTML BCR；CSS padding 进 OOXML 内 margin。
    # 这样：text-align:center 时居中点真的对齐 HTML 元素中心；
    #      padding-left 给 ::before marker 留位的设计在 PPT 也保持文字位置。
    style_padding = rec.get("style", {})
    tf.margin_left = px_to_emu(style_padding.get("paddingLeft", 0) or 0, SIZE_SCALE)
    tf.margin_right = px_to_emu(style_padding.get("paddingRight", 0) or 0, SIZE_SCALE)
    tf.margin_top = px_to_emu(style_padding.get("paddingTop", 0) or 0, SIZE_SCALE)
    tf.margin_bottom = px_to_emu(style_padding.get("paddingBottom", 0) or 0, SIZE_SCALE)
    # OOXML wrap 属性：square = 框内 wrap，none = 不 wrap 允许溢出
    # 必须与 tf.word_wrap 保持一致，否则 wrap 属性会覆盖 word_wrap 设置
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("wrap", "none" if no_auto_wrap else "square")
    # 竖排：CSS writing-mode → OOXML bodyPr vert
    # eaVert = 东亚竖排（CJK 字符保持正向，从右向左排列），覆盖 vertical-rl 主用例。
    # vertical-lr 用 eaVert 亦可（OOXML 没有 LTR 竖排原语，渲染器一般按列折行）。
    # sideways-* 在浏览器把字符整体旋转 90°，OOXML 没对应原语，落到 eaVert 兜底
    # 视觉会和浏览器有差异；纯 Latin 竖排建议用 transform: rotate(-90deg) 代替写法。
    wm = (rec.get("style", {}).get("writingMode") or "").lower()
    if wm and wm not in ("horizontal-tb", "lr", "lr-tb"):
        bodyPr.set("vert", "eaVert")
    # 垂直 anchor：HTML 用 flex/grid + align-items 居中时，OOXML 用 anchor 翻译
    style_for_anchor = rec.get("style", {})
    display = (style_for_anchor.get("display") or "").lower()
    align_items = (style_for_anchor.get("alignItems") or "").lower()
    anchor = "t"
    if "flex" in display or "grid" in display:
        if align_items == "center":
            anchor = "ctr"
        elif align_items in ("flex-end", "end"):
            anchor = "b"
    bodyPr.set("anchor", anchor)
    # 不要 autofit（防止 PPT 自己缩字）
    for child in list(bodyPr):
        if child.tag.endswith("normAutofit") or child.tag.endswith("spAutoFit"):
            bodyPr.remove(child)

    # 处理 runs
    runs = rec.get("runs", [])
    text_transform = rec.get("style", {}).get("textTransform", "none")

    # 第一段
    p = tf.paragraphs[0]
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    # 水平对齐：HTML flex/grid 容器靠 justify-content 居中，普通块靠 text-align。
    # OOXML 段落 algn 二选一：flex/grid 时 justify-content 优先（覆盖 text-align）。
    style_for_algn = rec.get("style", {})
    align = style_for_algn.get("textAlign", "start")
    align_map = {"start": "l", "left": "l", "center": "ctr", "right": "r", "end": "r"}
    justify_map = {"center": "ctr", "flex-end": "r", "end": "r", "right": "r",
                   "flex-start": "l", "start": "l", "left": "l",
                   "space-between": "just", "space-around": "ctr", "space-evenly": "ctr"}
    h_algn = align_map.get(align, "l")
    if "flex" in display or "grid" in display:
        jc = (style_for_algn.get("justifyContent") or "").lower()
        if jc in justify_map:
            h_algn = justify_map[jc]
    pPr.set("algn", h_algn)
    # 显式写行距：根据 HTML 实测 line-height/font-size 比率
    # 不写的话 PPT 用字体默认（CJK 字体默认行距大幅大于 CSS），导致大标题叠压下面元素
    # effective_lh_px / effective_fs 已在 no_auto_wrap 段计算（leaf vs runs 内最大字号 lh）
    _apply_line_spacing(
        pPr,
        f"{effective_lh_px}px" if effective_lh_px > 0 else rec.get("style", {}).get("lineHeight"),
        effective_fs,
    )

    # 清掉默认 run
    for r_el in p._p.findall(qn("a:r")):
        p._p.remove(r_el)

    first_para = p
    cur_para = first_para

    deduped = _normalize_runs(runs)

    pending_space = False
    for item in deduped:
        kind = item[0]
        if kind == "BREAK":
            _emit_line_break(cur_para, style_for_lh.get("fontSize", 16))
            pending_space = False
        elif kind == "SPACE":
            pending_space = True
        else:  # RUN
            _, run, body = item
            if pending_space:
                body = " " + body
                pending_space = False
            _emit_run(cur_para, body, run, text_transform)

    _fix_empty_paragraph_sizes(tf, style_for_lh.get("fontSize", 16))

    return tb


def _line_height_px(line_height, font_size_px: float) -> float:
    """CSS line-height → 像素值（用于估算文本框能容纳几条可视行）。"""
    if not font_size_px or font_size_px <= 0:
        return 0.0
    s = str(line_height or "").strip()
    if not s or s.lower() == "normal":
        return font_size_px * 1.2
    try:
        if s.endswith("px"):
            return float(s[:-2])
        if s.endswith("%"):
            return float(s[:-1]) / 100.0 * font_size_px
        return float(s) * font_size_px
    except ValueError:
        return font_size_px * 1.2


def _parse_line_spacing(line_height: str | float | None, font_size_px: float) -> tuple[str, int] | None:
    """Parse CSS line-height into DrawingML line spacing.

    CSS px line-height is an exact baseline distance, so emit spcPts instead
    of spcPct. PowerPoint/WPS interpret spcPct against the font's internal line
    box, which is usually taller than CSS line-height and looks too loose.

    Returns ("pts", value) for <a:spcPts> where value is 1/100 pt, or
    ("pct", value) for <a:spcPct> where value is 1/1000 percent.
    """
    if not line_height or not font_size_px or font_size_px <= 0:
        return None
    s = str(line_height).strip()
    if not s or s.lower() == "normal":
        return None
    # 百分号
    if s.endswith("%"):
        try:
            return ("pct", int(float(s[:-1]) * 1000))
        except ValueError:
            return None
    # px 显式
    if s.endswith("px"):
        try:
            px = float(s[:-2])
            return ("pts", max(100, int(round(px * PX_TO_PT * 100))))
        except ValueError:
            return None
    # 无单位数字（CSS line-height 倍数）
    try:
        return ("pct", int(float(s) * 100000))
    except ValueError:
        return None


def _apply_line_spacing(pPr, line_height, font_size_px, default_pct=120000):
    """在 pPr 上写 <a:lnSpc>。

    PPT 对 CJK 字体的默认行距比 CSS 大很多（典型 1.8-2.0x），不显式写就会出现
    大标题撑下来覆盖下方元素的"叠压"。这里把 CSS 的 line-height 显式传给 OOXML。

    CSS px 行高写成 <a:spcPts> 固定点值；百分比 / 倍数行高写成 <a:spcPct>。
    line-height 是 "normal" 时用 default_pct（120%）兜底——比 PPT 默认紧但
    不会产生叠压。
    """
    spacing = _parse_line_spacing(line_height, font_size_px)
    if spacing is None:
        spacing = ("pct", default_pct)
    # 清掉旧的 lnSpc（防重复运行）
    for old in pPr.findall(qn("a:lnSpc")):
        pPr.remove(old)
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    kind, val = spacing
    if kind == "pts":
        spc = etree.SubElement(lnSpc, qn("a:spcPts"))
    else:
        spc = etree.SubElement(lnSpc, qn("a:spcPct"))
    spc.set("val", str(val))
    # OOXML 子元素顺序：lnSpc 必须放在 pPr 的开头（在 buNone / buChar 之前）
    # pPr 子元素顺序: lnSpc, spcBef, spcAft, ..., defRPr
    # 简单做法：把 lnSpc 移到最前
    pPr.remove(lnSpc)
    pPr.insert(0, lnSpc)


def _emit_run(paragraph, text, run, text_transform):
    """向 paragraph 写入一个富文本 run。直接操作 OOXML 以精确控制 letterSpacing。"""
    if not text:
        return
    # text-transform: run 自带的 textTransform 优先（处理 descendant 上的 override，
    # 例如父级 text-transform: uppercase + .it-emph 子级 text-transform: none）
    effective_transform = run.get("textTransform") or text_transform
    if effective_transform == "uppercase":
        text = text.upper()
    elif effective_transform == "lowercase":
        text = text.lower()

    # 解析参数
    weight = run.get("fontWeight", "400")
    italic = run.get("fontStyle", "normal") == "italic"
    font_name, exact_weight_face = first_font_for_run(
        run.get("fontFamily", DEFAULT_LATIN_FALLBACK),
        weight,
        run.get("fontStyle", "normal"),
    )
    font_size_px = run.get("fontSize", 16)
    font_size_pt = round(font_size_px * PX_TO_PT, 2)

    try:
        bold = int(weight) >= 600
    except ValueError:
        bold = weight in ("bold", "bolder")
    if exact_weight_face:
        # The embedded typeface already carries the requested source weight and
        # style. Setting b/i on top can make PowerPoint synthesize another face.
        bold = False
        italic = False
    color_rgb = parse_rgb(run.get("color", "rgb(0,0,0)"))

    # letter-spacing: 字符串如 "3.6px" / "normal"
    ls = run.get("letterSpacing", "normal")
    if isinstance(ls, str) and ls.endswith("px"):
        ls_px = float(ls[:-2])
    else:
        ls_px = 0.0
    # OOXML spc 单位 = 1/100 pt
    spc_units = int(round(ls_px * PX_TO_PT * 100))

    # 中文 vs 英文：用 ea / latin 区分（CJK 范围定义见 text_utils.CJK_RE）
    is_chinese = is_cjk_text(text)

    # 构造 <a:r><a:rPr ...><a:rFont/></a:rPr><a:t>...</a:t></a:r>
    r_el = etree.SubElement(paragraph._p, qn("a:r"))
    rPr = etree.SubElement(r_el, qn("a:rPr"))
    rPr.set("lang", "zh-CN" if is_chinese else "en-US")
    rPr.set("sz", str(int(round(font_size_pt * 100))))  # OOXML sz = 1/100 pt
    if bold:
        rPr.set("b", "1")
    if italic:
        rPr.set("i", "1")
    if spc_units:
        rPr.set("spc", str(spc_units))
    rPr.set("dirty", "0")

    # fill color
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", "{:02X}{:02X}{:02X}".format(*color_rgb))

    # text-shadow → OOXML outerShdw（必须在 solidFill 之后、latin 之前）
    shadow = parse_text_shadow(run.get("textShadow", "none"))
    if shadow:
        import math
        dx, dy, blur, (sr, sg, sb, sa) = shadow
        dist_px = math.sqrt(dx * dx + dy * dy)
        if dist_px > 0.5 or blur > 0.5:
            angle_deg = math.degrees(math.atan2(dy, dx))
            if angle_deg < 0:
                angle_deg += 360
            effect_lst = etree.SubElement(rPr, qn("a:effectLst"))
            outer = etree.SubElement(effect_lst, qn("a:outerShdw"))
            outer.set("blurRad", str(int(blur * PX_TO_EMU * SIZE_SCALE)))
            outer.set("dist", str(int(dist_px * PX_TO_EMU * SIZE_SCALE)))
            outer.set("dir", str(int(angle_deg * 60000)))
            outer.set("algn", "ctr")
            outer.set("rotWithShape", "0")
            shd_clr = etree.SubElement(outer, qn("a:srgbClr"))
            shd_clr.set("val", f"{sr:02X}{sg:02X}{sb:02X}")
            if sa < 1.0:
                shd_a = etree.SubElement(shd_clr, qn("a:alpha"))
                shd_a.set("val", str(int(sa * 100000)))

    # font 分离：latin 用 css 第一项；ea 走 CJK 字体（Noto Serif/Sans SC）
    # PPT 会按字符自动用 latin 还是 ea，所以 IBM Plex Mono 文本里的 CJK 字符
    # 自动落到 Noto Sans SC 上，避免变 tofu。
    #
    # 关键：**只在 run 文本含 CJK 字符时**才写 <a:ea>。
    # PowerPoint 行为：ea 字体只对 CJK 字符生效，Latin 字符走 latin slot。
    # WPS 行为：会把 ea 字体的 advance width 应用到整 run（包括 Latin 字符）→
    # Noto Sans SC 是等宽 CJK（每字 ~1em），Latin 字母被撑到全角宽，文本横向溢出 +
    # 字间距异常宽。纯 Latin run 不写 ea，WPS 走 latin 字体的真实 advance。
    latin_el = etree.SubElement(rPr, qn("a:latin"))
    latin_el.set("typeface", font_name)
    if is_chinese:
        ea_name = cjk_font(run.get("fontFamily", ""), font_name)
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", ea_name)
    cs_el = etree.SubElement(rPr, qn("a:cs"))
    cs_el.set("typeface", font_name)

    t_el = etree.SubElement(r_el, qn("a:t"))
    t_el.text = text


def add_svg_picture(slide, rec):
    """直接用 measure 阶段已经截好的 SVG PNG。"""
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    png_path = rec.get("screenshot")
    if not png_path or not Path(png_path).exists():
        print(f"  [skip svg] no screenshot for {rec.get('marker', '?')}")
        return
    x, y, w, h = _resolve_rect(rec)
    slide.shapes.add_picture(png_path, x, y, w, h)


def add_img_picture(slide, rec):
    """measure 阶段截好的 <img> 元素 PNG（src 可能是 PNG/JPG/SVG/远程 URL）。
    走截图通道而不是直接嵌 src，避免 cross-origin / SVG 不能直嵌进 OOXML 等问题。"""
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    png_path = rec.get("screenshot")
    if not png_path or not Path(png_path).exists():
        print(f"  [skip img] no screenshot for {rec.get('src', '?')}")
        return
    x, y, w, h = _resolve_rect(rec)
    slide.shapes.add_picture(png_path, x, y, w, h)


def add_canvas_picture(slide, rec):
    """canvas 元素（Chart.js / WebGL / 自绘图）→ picture 嵌入。
    measure 阶段在切到目标页后等待 canvas 像素稳定再截图。
    """
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    png_path = rec.get("screenshot")
    if not png_path or not Path(png_path).exists():
        print(f"  [skip canvas] no screenshot for {rec.get('marker', '?')}")
        return
    x, y, w, h = _resolve_rect(rec)
    pic = slide.shapes.add_picture(png_path, x, y, w, h)
    _apply_rotation(pic, rec)
    return pic


def add_deco_snapshot(slide, rec):
    """装饰元素截图（background-image / box-shadow / 伪元素装饰 / 非平移 transform）
    → picture 嵌入。

    重要：Playwright `locator.screenshot()` 对带 transform 的元素截的是 AABB
    （旋转/skew 后的可见矩形），旋转已经"烘焙"到位图里。直接放在 AABB rect
    (rec.rect.x/y/w/h) 即可——**禁止调 `_apply_rotation`**：会双重旋转
    + naturalSize 压缩，对大尺寸旋转矩形（全宽 ribbon 等）尤其灾难。

    子节点的文字 / 子装饰按原流程绘制在它之上。
    """
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    png_path = rec.get("screenshot")
    if not png_path or not Path(png_path).exists():
        return
    x, y, w, h = _resolve_rect(rec)
    pic = slide.shapes.add_picture(png_path, x, y, w, h)
    return pic


def assemble_slide(slide, data):
    """装配一张 slide。"""
    bg_rgb = parse_rgb(data["slide"]["background"])
    bg_image = data["slide"].get("backgroundImage", "")
    add_background(slide, bg_rgb, bg_image)
    has_native_gradient = bool(
        bg_image and bg_image != "none" and "gradient" in bg_image.lower()
    )
    _prepare_text_layouts(data["records"])

    # ── Layout Agent 调优（opt-in）───────────────────────────────────
    # 设置环境变量 SLIDEFORGE_LAYOUT_AGENT=1 启用 LLM 智能布局调整。
    # 默认禁用：LLM 布局调整仍为实验性功能，在复杂内容页上可能产生位置偏差。
    import os
    if os.environ.get("SLIDEFORGE_LAYOUT_AGENT") == "1":
        try:
            from slideforge.agents.layout_agent import run_layout_agent
            from langchain_openai import ChatOpenAI

            llm = ChatOpenAI(model="gpt-4o", temperature=0)
            adjustments = run_layout_agent(llm, data["records"])
            for rec in data["records"]:
                eid = str(rec.get("id", ""))
                if eid in adjustments:
                    x, y, w, h = adjustments[eid]
                    rec["_adjusted_rect"] = (x, y, w, h)
        except ImportError:
            pass  # Layout Agent 模块未安装，使用 _scaled_rect
        except Exception:
            pass  # 任何异常 fallback 到 _scaled_rect

    text_records = []
    for rec in data["records"]:
        if rec["kind"] == "shape":
            # 整页 section 跳过（背景已由 add_background 铺）
            if rec["rect"]["w"] >= SLIDE_W_PX * 0.99 and rec["rect"]["h"] >= SLIDE_H_PX * 0.99:
                continue
            add_shape_box(slide, rec)
        elif rec["kind"] == "text":
            text_records.append(rec)
        elif rec["kind"] == "svg":
            add_svg_picture(slide, rec)
        elif rec["kind"] == "canvas":
            add_canvas_picture(slide, rec)
        elif rec["kind"] == "deco_snapshot":
            # 已用 native OOXML 渐变铺满背景时，跳过全屏截图
            if has_native_gradient and rec["rect"]["w"] >= SLIDE_W_PX * 0.99 and rec["rect"]["h"] >= SLIDE_H_PX * 0.99:
                continue
            add_deco_snapshot(slide, rec)
        elif rec["kind"] == "img":
            add_img_picture(slide, rec)

    # Text sits above rasterized SVG/canvas/deco snapshots. Otherwise an opaque
    # picture can cover positioned labels that belong visually on top of it.
    for rec in text_records:
        add_text_box(slide, rec)


def assemble(measurement, out_path: Path):
    """measurement 可以是 dict（in-process 调用）或 Path（CLI 调用）。"""
    if isinstance(measurement, (str, Path)):
        data = json.loads(Path(measurement).read_text(encoding="utf-8"))
    else:
        data = measurement

    # 兼容单页和多页
    if "slides" in data:
        slides_data = data["slides"]
    else:
        slides_data = [data]

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = _find_blank_layout(prs)

    for i, sdata in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        assemble_slide(slide, sdata)
        print(f"  page {i+1:02d}: {len(sdata.get('records', []))} records, theme={sdata['slide']['theme']}")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"saved {out_path} ({out_path.stat().st_size:,} B, {len(slides_data)} slides)")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    assemble(Path(sys.argv[1]).resolve(), Path(sys.argv[2]).resolve())
