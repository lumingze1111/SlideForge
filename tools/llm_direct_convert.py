#!/usr/bin/env python3
"""
让 DeepSeek LLM 直接根据 HTML 渲染 PPTX。

LLM 自己分析 HTML 结构，自己写 python-pptx 代码来生成 PPTX。
不调用项目已有的转换函数，由 LLM 全权负责渲染。

用法:
    python3 tools/llm_direct_convert.py
    python3 tools/llm_direct_convert.py --html path/to/slides.html
"""

import os
import sys
import re
import subprocess
import argparse
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent

SYSTEM_PROMPT = """你是一个 PPT 生成专家。你的任务是根据用户提供的 HTML 幻灯片文件，
用 python-pptx 库编写 Python 代码，生成视觉上高度还原的 PPTX 文件。

## 核心要求
1. **提取 HTML 中 ALL `.slide` 元素的全部内容** —— 标题、副标题、列表项、数据、装饰线、页码，一个都不能少
2. **忠实还原视觉风格**：background（渐变/纯色）、font-size、color、padding、text-align
3. 输出 **纯 Python 代码**（```python ... ``` 包裹），不输出任何解释文字
4. 代码必须可直接执行，包含完整 import 和 prs.save()

## 技术规则
- 幻灯片尺寸: width=Inches(13.33), height=Inches(7.5)（1280×720 比例）
- 渐变背景只使用 python-pptx 公开 API 能稳定表达的两端渐变；多段 CSS 渐变请降级为首尾两色或代表性纯色
- 文字用 paragraph.font.size / color.rgb / bold / alignment
- 文本框大小和位置根据 HTML 布局合理估算
- 用 MSO_SHAPE.RECTANGLE 实现装饰线和色块
- HTML 中的 div 网格布局，用多个 textbox 或 shape 模拟
- 字体默认使用 'PingFang SC' 或 'Microsoft YaHei'

## python-pptx 常见陷阱（必须遵守！）
1. **渐变背景保守处理**: 只设置 `fill.gradient()` 默认已有的两个 stop。不要扩展 stop 数量，也不要访问 `_fill`、`_GradFill`、`_element` 等私有对象。
   ```python
   fill.gradient()
   stops = fill.gradient_stops
   stops[0].position = 0.0
   stops[0].color.rgb = color1
   stops[1].position = 1.0
   stops[1].color.rgb = color2
   ```
2. **slide.background.fill 与 shape.fill 不同**: slide 背景的 fill 对象 API 与形状的 fill 相同，直接用即可
3. **text_frame 首个段落**: `tf.paragraphs[0]` 始终存在，无需手动创建
4. **颜色值**: RGBColor 参数接受 0-255 整数，16进制写法 `RGBColor(0xFF, 0xFF, 0xFF)`
5. **导入必须完整**: `from pptx import Presentation; from pptx.util import Inches, Pt; from pptx.dml.color import RGBColor; from pptx.enum.text import PP_ALIGN`

## 内容完整性检查（非常重要！）
HTML 中每一页（.slide）的 **全部文字内容** 都必须出现在 PPTX 中。
- h1, h2 标题 → 幻灯片标题
- p 段落 → 文本框
- li 列表 → 多行文本框
- div 中的文字块 → 对应位置的文本框
- 所有 bullet points、数据、引用、说明文字都要保留

现在开始。"""


def read_html(html_path: Path) -> str:
    with open(html_path, "r", encoding="utf-8") as f:
        return f.read()


def strip_base64_for_llm(html: str) -> tuple[str, list[dict]]:
    """
    从 HTML 中剥离 base64 图片数据，用占位符替换。
    返回 (精简后的 HTML, 图片信息列表)

    这样 LLM 只需处理文本和布局，图片后续单独插入。
    """
    images = []
    counter = [0]

    def replace_base64(match):
        full_tag = match.group(0)
        counter[0] += 1
        img_id = f"IMG_PLACEHOLDER_{counter[0]}"

        # 提取 style 中的位置信息
        style_match = re.search(r'style="([^"]*)"', full_tag)
        style = style_match.group(1) if style_match else ""

        # 提取 alt 信息
        alt_match = re.search(r'alt="([^"]*)"', full_tag)
        alt = alt_match.group(1) if alt_match else ""

        images.append({
            "id": img_id,
            "style": style,
            "alt": alt,
            "index": counter[0],
        })

        return f'<img src="[{img_id}]" style="{style}" alt="{alt}">'

    # 匹配包含 base64 数据的 img 标签
    pattern = r'<img\s+[^>]*src="data:image/[^"]*"[^>]*>'
    stripped = re.sub(pattern, replace_base64, html)

    return stripped, images


def truncate_html_for_llm(html: str, max_chars: int = 50000) -> str:
    """
    如果 HTML 仍然过长，截取前 N 个字符并标注。
    """
    if len(html) <= max_chars:
        return html
    return html[:max_chars] + "\n\n<!-- HTML 已截断，后续内容省略 -->"


def call_deepseek(user_prompt: str, api_key: str) -> str:
    """调用 DeepSeek API，返回生成的 Python 代码"""
    from openai import OpenAI
    client = OpenAI(base_url="https://api.deepseek.com", api_key=api_key)

    resp = client.chat.completions.create(
        model="deepseek-chat",
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.1,
        max_tokens=8192,
    )
    return resp.choices[0].message.content


def extract_code(text: str) -> str:
    """从 LLM 回复中提取 Python 代码"""
    pattern = r"```python\s*\n(.*?)```"
    match = re.search(pattern, text, re.DOTALL)
    if match:
        return match.group(1).strip()
    # fallback: 匹配任意代码块
    pattern = r"```\s*\n(.*?)```"
    match = re.search(pattern, text, re.DOTALL)
    if match:
        return match.group(1).strip()
    return text.strip()


def extract_notes(html_path: Path) -> dict[int, str]:
    """从 HTML 的 data-notes 属性提取演讲者备注。"""
    html = html_path.read_text(encoding="utf-8")
    pattern = re.compile(
        r'<div[^>]*?data-pptx-slide[^>]*?data-notes\s*=\s*"([^"]*)"',
        re.IGNORECASE,
    )
    notes_map = {}
    for i, m in enumerate(pattern.finditer(html), start=1):
        notes = m.group(1)
        notes = notes.replace("&#39;", "'").replace("&quot;", '"').replace("&amp;", "&")
        notes = notes.replace("&lt;", "<").replace("&gt;", ">")
        notes = notes.replace("\\n", "\n")
        if notes.strip():
            notes_map[i] = notes
    return notes_map


def inject_notes(pptx_path: Path, notes_map: dict[int, str]) -> int:
    """向 PPTX 注入演讲者备注，返回注入页数。"""
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    injected = 0
    for idx, slide in enumerate(prs.slides, start=1):
        notes_text = notes_map.get(idx)
        if not notes_text:
            continue
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.clear()
            tf.paragraphs[0].text = notes_text
            injected += 1
        except Exception:
            pass
    if injected:
        prs.save(str(pptx_path))
    return injected


def inject_images_to_pptx(pptx_path: Path, html_content: str) -> int:
    """
    从原始 HTML 中提取 base64 图片，注入到 PPTX 的对应页面中。
    只处理背景图片和 inline 图片。
    """
    import base64
    import io
    from pptx import Presentation
    from pptx.util import Inches, Emu

    # 提取每个 slide 中的图片
    slide_pattern = re.compile(
        r'<div\s+class="slide"[^>]*>(.*?)</div>\s*(?=<div\s+class="slide"|</div>\s*</div>\s*</body>)',
        re.DOTALL
    )
    img_pattern = re.compile(
        r'<img\s+src="data:image/(jpeg|png);base64,([^"]+)"[^>]*>',
        re.DOTALL
    )

    # 按 slide 分组提取图片
    slide_images = {}
    for slide_idx, slide_match in enumerate(slide_pattern.finditer(html_content)):
        slide_html = slide_match.group(1)
        imgs = img_pattern.findall(slide_html)
        if imgs:
            slide_images[slide_idx] = imgs  # [(format, base64_data), ...]

    if not slide_images:
        return 0

    try:
        prs = Presentation(str(pptx_path))
        injected = 0

        for slide_idx, imgs in slide_images.items():
            if slide_idx >= len(prs.slides):
                continue

            slide = prs.slides[slide_idx]
            # 只插入第一张图片作为背景（避免过多图片）
            img_format, img_b64 = imgs[0]

            try:
                img_bytes = base64.b64decode(img_b64)
                img_stream = io.BytesIO(img_bytes)

                # 作为全页背景图片插入（放在最底层）
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                pic = slide.shapes.add_picture(
                    img_stream,
                    left=0, top=0,
                    width=slide_width,
                    height=slide_height
                )

                # 移动到最底层
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)

                # 设置透明度（通过 alpha 通道模拟）
                from pptx.oxml.ns import qn
                from lxml import etree
                blipFill = pic._element.find(qn('p:blipFill'))
                if blipFill is not None:
                    blip = blipFill.find(qn('a:blip'))
                    if blip is not None:
                        alphaModFix = etree.SubElement(blip, qn('a:alphaModFix'))
                        alphaModFix.set('amt', '25000')  # 25% 不透明度

                injected += 1

            except Exception:
                continue

        if injected:
            prs.save(str(pptx_path))

        return injected

    except Exception:
        return 0


def main():
    parser = argparse.ArgumentParser(description="DeepSeek LLM 直接渲染 HTML → PPTX")
    parser.add_argument("--html", default=str(PROJECT_ROOT / "output" / "slides_你好 旅行者.html"))
    parser.add_argument("--output", "-o", default=None)
    args = parser.parse_args()

    api_key = os.environ.get("DEEPSEEK_API_KEY")
    if not api_key:
        print("❌ 请设置环境变量 DEEPSEEK_API_KEY")
        print("   export DEEPSEEK_API_KEY='your-api-key-here'")
        sys.exit(1)

    html_path = Path(args.html).resolve()
    if not html_path.exists():
        print(f"❌ 文件不存在: {html_path}")
        sys.exit(1)

    if args.output:
        output_path = Path(args.output).resolve()
    else:
        stem = html_path.stem.replace(" ", "_")
        output_path = html_path.parent / f"{stem}_llm.pptx"

    print(f"📄 读取 HTML: {html_path}")
    html_content = read_html(html_path)
    print(f"   文件大小: {len(html_content):,} 字符")

    # 剥离 base64 图片数据，避免超出 API 上下文限制
    stripped_html, image_placeholders = strip_base64_for_llm(html_content)
    stripped_html = truncate_html_for_llm(stripped_html, max_chars=50000)
    print(f"   精简后大小: {len(stripped_html):,} 字符（剥离 {len(image_placeholders)} 张图片）")

    # 构建 prompt
    image_note = ""
    if image_placeholders:
        image_note = f"""
注意：HTML 中原有 {len(image_placeholders)} 张 base64 图片已用占位符 [IMG_PLACEHOLDER_N] 替换。
在生成的代码中，请用注释标记图片位置即可（如 # TODO: insert image here），不要尝试嵌入图片数据。
图片会在后续步骤单独处理。"""

    prompt = f"""请将以下 HTML 幻灯片文件完整转换为 PPTX。

HTML 文件: {html_path}
输出 PPTX: {output_path}
{image_note}

精简后的 HTML 内容如下：
```html
{stripped_html}
```

请逐页分析 HTML 中每一个 .slide 元素，提取全部文字内容，
然后用 python-pptx 生成 PPTX。确保每页的标题、列表、数据、装饰元素都完整保留。
对于图表（CSS 柱状图、饼图等），请用 python-pptx 的图表 API 或形状模拟还原。

注意：无需处理演讲者备注，后续会单独注入。"""

    print(f"\n🤖 调用 DeepSeek 生成 PPTX 代码...")
    reply = call_deepseek(prompt, api_key)
    print(f"   回复长度: {len(reply):,} 字符")

    code = extract_code(reply)
    code_path = output_path.with_suffix(".py")

    # 补全 output_path 变量
    if "output_path" not in code:
        code = code.replace("prs.save(", f"output_path = r'{output_path}'\nprs.save(")

    # 保存生成的代码
    code_path.write_text(code, encoding="utf-8")
    print(f"📝 代码已保存: {code_path}")

    max_fix_attempts = 3
    for attempt in range(max_fix_attempts + 1):
        print(f"\n⚡ 执行生成的 PPTX 代码 (第 {attempt + 1}/{max_fix_attempts + 1} 次)...")
        result = subprocess.run(
            [sys.executable, str(code_path)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode == 0:
            break
        print(f"❌ 执行失败:")
        print(result.stderr[:2000])
        if attempt < max_fix_attempts:
            print(f"\n🤖 将错误发给 DeepSeek 修复 (剩余 {max_fix_attempts - attempt} 次)...")
            fix_prompt = f"生成的代码执行报错:\n{result.stderr[:2000]}\n\n请修复代码:\n```python\n{code}\n```"
            fix_reply = call_deepseek(
                f"请修复以下代码错误。\nHTML 文件: {html_path}\n\n{fix_prompt}\n\n输出 PPTX 保存到: {output_path}",
                api_key
            )
            code = extract_code(fix_reply)
            code_path.write_text(code, encoding="utf-8")
            print(f"📝 已更新代码: {code_path}")
        else:
            print(f"❌ 已尝试 {max_fix_attempts} 次修复，仍然失败")
            sys.exit(1)

    if output_path.exists():
        # 注入图片（从原始 HTML 中提取 base64 图片并插入 PPTX）
        if image_placeholders:
            injected_imgs = inject_images_to_pptx(output_path, html_content)
            if injected_imgs:
                print(f"   🖼️ 已注入 {injected_imgs} 张图片")

        # 注入演讲者备注
        notes_map = extract_notes(html_path)
        if notes_map:
            injected = inject_notes(output_path, notes_map)
            print(f"   📝 已注入 {injected} 页演讲者备注")
        print(f"\n✅ PPTX 已生成: {output_path}")
        print(f"   文件大小: {output_path.stat().st_size / 1024:.1f} KB")
    else:
        print(f"\n⚠️ 文件未生成，检查代码 {code_path}")


if __name__ == "__main__":
    main()
